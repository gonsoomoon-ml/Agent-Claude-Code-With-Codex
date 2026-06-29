# ruff: noqa  — 일회성 슬라이드 생성기(레이아웃 좌표 밀집). src 린트 기준 비적용.
"""solution-architecture.pptx — 다크 테마 솔루션 아키텍처 슬라이드 생성기. 16:9, Noto Sans CJK KR.

재생성:  uv run --with python-pptx python design/architecture/make_slide.py
미리보기(PNG): libreoffice --headless --convert-to png --outdir <dir> design/architecture/solution-architecture.pptx
※ 출력 .pptx 는 .gitignore(바이너리) — 이 스크립트가 진실의 원천.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

FONT = "Noto Sans CJK KR"
BG = RGBColor(0x0D, 0x0D, 0x10)
TEXT = RGBColor(0xF2, 0xF2, 0xEF)
SUB = RGBColor(0xBC, 0xBC, 0xB8)
INKD = RGBColor(0x14, 0x0C, 0x08)        # on-light/on-coral 어두운 텍스트(대비↑)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CORAL = RGBColor(0xFF, 0x6B, 0x47)
CARD = RGBColor(0x1C, 0x1D, 0x22)
PANEL = RGBColor(0x14, 0x15, 0x19)
ARR = RGBColor(0x8A, 0x8A, 0x92)
SEP = RGBColor(0x33, 0x33, 0x3A)
EDGE = RGBColor(0x8F, 0xB4, 0xD6)
ACORE = RGBColor(0x3F, 0xD3, 0xBF)
COMP = RGBColor(0xE3, 0xB3, 0x4E)        # 앰버 — 코랄과 분리(P1-3)
DATA = RGBColor(0x5B, 0x9B, 0xD5)
DELIV = RGBColor(0x5B, 0xC9, 0x8C)
SRC = RGBColor(0x9A, 0xA0, 0xA6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG


def _font(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size); run.font.color.rgb = color; run.font.bold = bold
    run.font.italic = italic; run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", FONT)


def box(x, y, w, h, lines, border, fill, lw=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.CENTER, ml=3):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = border; sp.line.width = Pt(lw); sp.shadow.inherit = False
    try:
        sp.adjustments[0] = min(0.085 / min(w, h), 0.5)   # 곡률 절대값 통일(P2-6)
    except Exception:
        pass
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_top", "margin_bottom", "margin_right"):
        setattr(tf, m, Pt(3))
    tf.margin_left = Pt(ml)
    for i, (txt, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_before = Pt(0); p.space_after = Pt(0)
        r = p.add_run(); r.text = txt; _font(r, size, color, bold)
    return sp


def text(x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(1))
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.alignment = align; p.space_before = Pt(0); p.space_after = Pt(0)
        for tup in line:
            r = p.add_run(); r.text = tup[0]; _font(r, tup[1], tup[3], tup[2])
    return tb


def badge(cx, cy, n, d=0.32):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = CORAL; sp.line.color.rgb = BG; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(0))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n); _font(r, 11.5, INKD, True)   # 잉크 숫자(P1-1)
    return sp


def tab(bx, by, n):       # 박스 좌상단 *모서리* = 번호 탭(글자 안 가림)
    badge(bx, by, n)


def arrow(x1, y1, x2, y2, color=ARR, w=1.75):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w); cn.shadow.inherit = False
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return cn


def line(x1, y1, x2, y2, color=ARR, w=1.1):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w); cn.shadow.inherit = False
    return cn


# ── Title ──
text(0.45, 0.16, 12.4, 0.5, [[("Solution Architecture", 23, True, WHITE),
                              ("   verify-before-publish AI 브리핑", 15, False, SUB)]])
text(0.47, 0.64, 12.4, 0.3, [[("AWS · us-east-1   —   Claude(author) ⇄ Codex(certifier) 듀얼 하니스 · 검증 후 발행", 11, False, ACORE)]])
ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.45), Inches(0.97), Inches(12.88), Inches(0.97))
ln.line.color.rgb = CORAL; ln.line.width = Pt(2.5); ln.shadow.inherit = False

# ── (1) 진입 ──
text(0.45, 1.08, 4, 0.3, [[("진입 (trigger)", 13, True, EDGE)]])
ey = 1.48
box(0.55, ey, 1.25, 0.6, [("USER", 12, True, TEXT)], EDGE, CARD)
box(2.05, ey, 1.55, 0.6, [("CloudFront", 12, True, TEXT), ("+ S3 · SPA", 8.6, False, SUB)], EDGE, CARD)
box(3.85, ey, 1.6, 0.6, [("API Gateway", 12, True, TEXT), ("HTTP API + Lambda", 8.6, False, SUB)], EDGE, CARD)
box(5.7, ey, 1.45, 0.6, [("Cognito", 12, True, TEXT), ("인가 (PKCE/M2M)", 8.6, False, SUB)], EDGE, CARD)
box(8.35, ey, 1.7, 0.6, [("EventBridge", 12, True, TEXT), ("cron 07:00 KST", 8.6, False, SUB)], EDGE, CARD)
box(10.3, ey, 1.55, 0.6, [("Lambda", 12, True, TEXT), ("dispatch", 8.6, False, SUB)], EDGE, CARD)
for (a, b) in [(1.8, 2.05), (3.6, 3.85), (5.45, 5.7), (10.05, 10.3)]:
    arrow(a, ey + 0.3, b, ey + 0.3, ARR, 1.5)
text(2.05, ey + 0.63, 5.1, 0.25, [[("웹 — 체험 / 구독", 9, False, SUB)]], PP_ALIGN.CENTER)
text(8.35, ey + 0.63, 3.5, 0.25, [[("스케줄 — 매일 발송", 9, False, SUB)]], PP_ALIGN.CENTER)
tab(0.55, ey, 1)

# ── (2) invoke (웹·스케줄 두 경로가 합류) ──
line(4.65, ey + 0.6, 6.55, 2.14, EDGE, 1.4)        # 웹(API GW) ↘
line(11.07, ey + 0.6, 6.65, 2.14, EDGE, 1.4)       # 스케줄(Lambda) ↙
arrow(6.6, 2.14, 6.6, 2.66, CORAL, 2.5)
badge(6.6, 2.4, 2)
text(6.85, 2.24, 4.0, 0.3, [[("invoke (trial / scheduled)", 9.5, False, CORAL)]])

# ── AgentCore Runtime ──
rt_x, rt_y, rt_w, rt_h = 0.55, 2.66, 10.7, 1.55
box(rt_x, rt_y, rt_w, rt_h, [], ACORE, PANEL, lw=2.25)
text(rt_x + 0.2, rt_y + 0.09, 8, 0.3, [[("AgentCore Runtime", 14, True, ACORE),
                                        ("   (claude + codex 번들)", 9, False, SUB)]])
sy = rt_y + 0.52
box(0.85, sy, 1.65, 0.78, [("retrieval", 11.5, True, TEXT), ("출처 페치", 8.6, False, SUB)], ACORE, CARD); tab(0.85, sy, 3)
box(2.9, sy, 1.65, 0.78, [("동결 · 필터", 11.5, True, TEXT), ("sha256 · 관련성", 8.6, False, SUB)], ACORE, CARD); tab(2.9, sy, 4)
# GATE = 코랄 히어로 + author/certifier 서브카드 ((5)(6) 텍스트 인라인 — 글자 안 가림)
gx, gw, gh = 4.95, 4.1, 1.06
box(gx, sy - 0.08, gw, gh, [], CORAL, CORAL, lw=1.0)
text(gx, sy - 0.02, gw, 0.22, [[("GATE — 결정론 = 신뢰의 원천", 10, True, INKD)]], PP_ALIGN.CENTER)
box(gx + 0.16, sy + 0.22, 1.78, 0.46, [("(5) author · Claude", 9, True, INKD)], WHITE, WHITE, lw=0.75)
text(gx + 1.97, sy + 0.27, 0.34, 0.3, [[("⇄", 14, True, INKD)]], PP_ALIGN.CENTER)
box(gx + 2.32, sy + 0.22, 1.62, 0.46, [("(6) certifier · Codex", 9, True, INKD)], WHITE, WHITE, lw=0.75)
text(gx, sy + 0.72, gw, 0.22, [[("★ decorrelation · envelope 4필드", 8.5, True, INKD)]], PP_ALIGN.CENTER)
box(9.25, sy, 1.65, 0.78, [("render", 11.5, True, TEXT), ("검증 명세서", 8.6, False, SUB)], ACORE, CARD); tab(9.25, sy, 7)
arrow(2.5, sy + 0.39, 2.9, sy + 0.39, ARR, 1.75)
arrow(4.55, sy + 0.39, 4.95, sy + 0.4, ARR, 1.75)
arrow(9.05, sy + 0.39, 9.25, sy + 0.39, ARR, 1.75)

# ── 백킹 서비스 ──
text(0.45, 4.3, 4, 0.25, [[("백킹 서비스 (backing)", 13, True, EDGE)]])
by = 4.66
box(0.85, by, 1.85, 0.72, [("AgentCore Gateway", 11, True, TEXT), ("retrieval MCP · 기본 off", 8.6, False, SUB)], ACORE, CARD)
box(2.95, by, 1.7, 0.72, [("AgentCore Identity", 11, True, TEXT), ("토큰 · 볼트", 8.6, False, SUB)], ACORE, CARD)
box(5.3, by, 2.3, 0.72, [("Bedrock", 11.5, True, TEXT), ("Claude + GPT-5.5 (cross-family)", 8.6, False, SUB)], COMP, CARD)
box(7.85, by, 1.9, 0.72, [("DynamoDB", 11.5, True, TEXT), ("store · cache · ledger · users", 8.4, False, SUB)], DATA, CARD)
box(10.0, by, 1.85, 0.72, [("SES", 11.5, True, TEXT), ("이메일 발송", 8.6, False, SUB)], DELIV, CARD)
arrow(1.75, rt_y + rt_h, 1.75, by, ACORE, 1.5)         # retrieval cluster
arrow(6.45, rt_y + rt_h, 6.45, by, COMP, 1.5)          # GATE → Bedrock
arrow(8.8, rt_y + rt_h, 8.8, by, DATA, 1.5)            # → DynamoDB (store)
arrow(10.2, rt_y + rt_h, 10.7, by, DELIV, 1.75)        # render → SES (send)
badge(10.45, (rt_y + rt_h + by) / 2, 8)
line(2.95, by + 0.36, 2.72, by + 0.36, ACORE, 1.2)     # Identity → Gateway
text(2.72, by + 0.72, 2.0, 0.22, [[("token · Cognito 인가", 8, False, ACORE)]], PP_ALIGN.CENTER)

# ── Sources / USER inbox ──
sy2 = 5.72
box(0.85, sy2, 1.85, 0.6, [("Sources (RSS)", 10.5, True, TEXT), ("catalog.yaml", 8.4, False, SUB)], SRC, CARD)
arrow(1.75, by + 0.72, 1.75, sy2, SRC, 1.5)
box(10.0, sy2, 1.85, 0.6, [("USER 받은편지함", 10.5, True, TEXT), ("매일 07:00 KST", 8.4, False, SUB)], DELIV, CARD)
arrow(10.9, by + 0.72, 10.9, sy2, DELIV, 1.75)
tab(10.0, sy2, 9)

# ── 범례 ──
lgy = 6.5
line(0.45, lgy - 0.1, 12.88, lgy - 0.1, SEP, 1.0)
left = [
    "(1) 진입 — USER→CloudFront→API Gateway(Cognito 인가) · EventBridge→Lambda",
    "(2) AgentCore Runtime invoke (체험 / 스케줄)",
    "(3) retrieval — Runtime→Sources 페치 (Gateway MCP 경유는 옵션·기본 off · token=Identity)",
    "(4) 동결(sha256) + 관련성(토픽) 필터",
    "(5) author = Claude (Bedrock)",
]
right = [
    "(6) certifier = GPT-5.5 (Bedrock) — ★ decorrelation (envelope 4필드)",
    "(7) render (검증 명세서)",
    "(8) DynamoDB 저장·캐시  +  SES 발송",
    "(9) USER 받은편지함 (매일 07:00 KST)",
    "핵심 — certifier 는 envelope 4필드만, 작성자 추론 미전달 → 독립 재도출",
]
text(0.45, lgy, 6.55, 1.0, [[(s, 9.3, False, SUB)] for s in left])
text(6.95, lgy, 6.0, 1.0, [[(s, 9.3, (i == 4), CORAL if i == 4 else SUB)] for i, s in enumerate(right)])

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "solution-architecture.pptx")
prs.save(out)
print("saved", out)
